"""Index local documents (PDF, Excel, CSV, code/text) and answer
questions about them. Everything stays on disk in ChromaDB."""

from pathlib import Path

import chromadb
import ollama
import pandas as pd
from pypdf import PdfReader

from app.config import (
    CHUNK_OVERLAP,
    CHUNK_SIZE,
    CODE_EXTENSIONS,
    DOCUMENTS_DIR,
    EMBED_MODEL,
    RERANK_CANDIDATES,
    RERANKER_MODEL,
    TOP_K,
    VECTOR_DB_DIR,
)
from app.modelstate import current_model

_client = chromadb.PersistentClient(path=str(VECTOR_DB_DIR))
COLLECTION_NAME = "documents"

# noise directories inside cloned repos — never worth indexing
SKIP_DIRS = {
    ".git", ".hg", ".svn", "node_modules", "__pycache__",
    ".venv", "venv", "dist", "build", ".idea", ".vscode",
}
MAX_FILE_BYTES = 1_000_000  # skip huge files (minified bundles, data dumps)

_reranker = None


def _rerank(question: str, docs: list[str], sources: list[str]):
    """Re-score candidate chunks with a cross-encoder so only the most
    relevant reach the LLM. Falls back to embedding order if the
    reranker can't load."""
    global _reranker
    try:
        if _reranker is None:
            from sentence_transformers import CrossEncoder

            _reranker = CrossEncoder(RERANKER_MODEL)
        scores = _reranker.predict([(question, doc) for doc in docs])
        order = sorted(range(len(docs)), key=lambda i: scores[i], reverse=True)
        ranked = order[:TOP_K]
        return [docs[i] for i in ranked], [sources[i] for i in ranked]
    except Exception as exc:
        print(f"[rag] reranker unavailable, using embedding order: {exc}")
        return docs[:TOP_K], sources[:TOP_K]


def _read_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _read_spreadsheet(path: Path) -> str:
    if path.suffix.lower() == ".csv":
        sheets = {"csv": pd.read_csv(path)}
    else:
        sheets = pd.read_excel(path, sheet_name=None)
    parts = []
    for name, df in sheets.items():
        parts.append(f"## Sheet: {name}\n{df.to_string(max_rows=500)}")
    return "\n\n".join(parts)


def _read_word(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _read_powerpoint(path: Path) -> str:
    from pptx import Presentation

    parts = []
    for i, slide in enumerate(Presentation(str(path)).slides, 1):
        texts = [
            shape.text for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        parts.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _read_file(path: Path) -> str | None:
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _read_pdf(path)
        if suffix in (".xlsx", ".xls", ".csv"):
            return _read_spreadsheet(path)
        if suffix == ".docx":
            return _read_word(path)
        if suffix == ".pptx":
            return _read_powerpoint(path)
        if suffix in CODE_EXTENSIONS:
            return path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:  # skip unreadable files, keep indexing the rest
        print(f"[rag] could not read {path.name}: {exc}")
    return None


def _chunk(text: str) -> list[str]:
    chunks = []
    step = CHUNK_SIZE - CHUNK_OVERLAP
    for start in range(0, len(text), step):
        piece = text[start : start + CHUNK_SIZE].strip()
        if piece:
            chunks.append(piece)
    return chunks


def _embed(texts: list[str]) -> list[list[float]]:
    return ollama.embed(model=EMBED_MODEL, input=texts)["embeddings"]


def index_documents() -> str:
    """(Re)build the index from everything in data/documents/."""
    try:
        _client.delete_collection(COLLECTION_NAME)
    except Exception:
        pass
    collection = _client.create_collection(COLLECTION_NAME)

    files = [
        p
        for p in sorted(DOCUMENTS_DIR.rglob("*"))
        if p.is_file()
        and not (set(p.parts) & SKIP_DIRS)
        and p.stat().st_size <= MAX_FILE_BYTES
    ]
    if not files:
        return f"No files found. Put PDFs, Excel files, or code into {DOCUMENTS_DIR}"

    indexed, skipped = 0, 0
    for path in files:
        text = _read_file(path)
        if not text or not text.strip():
            skipped += 1
            continue
        chunks = _chunk(text)
        rel = str(path.relative_to(DOCUMENTS_DIR))
        # embed in batches to keep requests small
        for batch_start in range(0, len(chunks), 32):
            batch = chunks[batch_start : batch_start + 32]
            collection.add(
                ids=[f"{rel}::{batch_start + i}" for i in range(len(batch))],
                documents=batch,
                embeddings=_embed(batch),
                metadatas=[{"source": rel}] * len(batch),
            )
        indexed += 1

    return f"Indexed {indexed} file(s) ({skipped} skipped) from {DOCUMENTS_DIR}"


def retrieve_context(question: str) -> tuple[str, list[str]] | None:
    """Best document chunks for a question as (context, sources),
    or None if there is no index / nothing relevant."""
    try:
        collection = _client.get_collection(COLLECTION_NAME)
    except Exception:
        return None

    result = collection.query(
        query_embeddings=_embed([question]),
        n_results=min(RERANK_CANDIDATES, collection.count()),
    )
    docs = result["documents"][0]
    sources = [m["source"] for m in result["metadatas"][0]]
    if not docs:
        return None
    docs, sources = _rerank(question, docs, sources)
    context = "\n\n---\n\n".join(
        f"[{src}]\n{doc}" for src, doc in zip(sources, docs)
    )
    return context, sources


def ask_documents(question: str) -> str:
    """Retrieve relevant chunks and answer with the local LLM."""
    retrieved = retrieve_context(question)
    if retrieved is None:
        return (
            "Nothing relevant found — is the index built? "
            "Click 'Index documents' first."
        )
    context, sources = retrieved
    response = ollama.chat(
        model=current_model(),
        messages=[
            {
                "role": "system",
                "content": (
                    "Answer using ONLY the provided document excerpts. "
                    "Cite the source file names you used. If the answer "
                    "is not in the excerpts, say so."
                ),
            },
            {
                "role": "user",
                "content": f"Documents:\n{context}\n\nQuestion: {question}",
            },
        ],
    )
    answer = response["message"]["content"]
    unique_sources = ", ".join(dict.fromkeys(sources))
    return f"{answer}\n\n_Sources searched: {unique_sources}_"
