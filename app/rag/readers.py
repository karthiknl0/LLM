"""Document readers for the RAG pipeline: PDF, spreadsheet, Word,
PowerPoint, and plain text / code files."""

from pathlib import Path

from app.core.config import CODE_EXTENSIONS


def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _read_spreadsheet(path: Path) -> str:
    import pandas as pd
    if path.suffix.lower() == ".csv":
        sheets = {"csv": pd.read_csv(path)}
    else:
        sheets = pd.read_excel(path, sheet_name=None)
    parts = []
    for name, df in sheets.items():
        parts.append(f"## Sheet: {name}\n{df.to_string(max_rows=500)}")
    return "\n\n".join(parts)


def _read_word(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _read_powerpoint(path: Path) -> str:
    from pptx import Presentation

    parts = []
    for i, slide in enumerate(Presentation(str(path)).slides, 1):
        texts = [
            shape.text for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        parts.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(parts)


def read_file(path: Path) -> str | None:
    """Read a document file; return text or None if unreadable/unsupported."""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _read_pdf(path)
        if suffix in (".xlsx", ".xls", ".csv"):
            return _read_spreadsheet(path)
        if suffix == ".docx":
            return _read_word(path)
        if suffix == ".pptx":
            return _read_powerpoint(path)
        if suffix in CODE_EXTENSIONS:
            return path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:  # skip unreadable files, keep indexing the rest
        print(f"[rag] could not read {path.name}: {exc}")
    return None
